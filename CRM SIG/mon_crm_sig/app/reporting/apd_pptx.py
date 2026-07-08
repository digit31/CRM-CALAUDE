"""
apd_pptx.py - Génération du livrable APD (plan PDF) par remplissage du template
PowerPoint « APD_HTL_NOM PROJET.pptx », puis export PPTX -> PDF.

Le template (4 slides A4) est copié puis rempli par projet :
  · Slide 1 : en-tête (code + type), adresse, grand plan général (image générée
              depuis les SHP, ou image uploadée), pied de page.
  · Slide 2 : « Caractéristiques de la liaison » (données Console Étude).
  · Slide 3 : plan de masse (image uploadée) — uniquement si site = BTS,
              sinon la slide est supprimée.
  · Slide 4 : une vignette par BPE À CRÉER (nom, adresse, chambre/appui, modèle).

Export PPTX -> PDF via PowerPoint COM (pywin32) — PowerPoint requis sur la machine.

Dépendances : python-pptx, geopandas ; pywin32 pour l'export.
"""

import os
import copy
import math
import logging

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

import geopandas as gpd

from app.reporting import apd_generator

logger = logging.getLogger("crm_sig.apd_pptx")

# Couleurs des appuis (identiques au livrable de référence)
_ROUGE = RGBColor(0xC0, 0x00, 0x00)
_ORANGE = RGBColor(0xE8, 0xB2, 0x1C)
_BLEU = RGBColor(0x2D, 0x9B, 0xF0)


# ---------------------------------------------------------------------------
# Helpers PPTX
# ---------------------------------------------------------------------------

def _tables(slide):
    return [sh for sh in slide.shapes if sh.has_table]


def _header_table(slide):
    for sh in _tables(slide):
        if sh.top is not None and sh.top < Inches(1.5) and len(sh.table.columns) == 3:
            return sh
    return None


def _footer_table(slide):
    for sh in _tables(slide):
        if sh.top is not None and sh.top > Inches(9.5) and len(sh.table.columns) == 4:
            return sh
    return None


def _base_rpr(text_frame):
    """Clone le format (rPr) du 1er run — préserve taille, gras ET couleur de
    thème (SCHEME) que l'API haut niveau ne sait pas recopier."""
    for p in text_frame.paragraphs:
        for r in p.runs:
            rpr = r._r.find(qn("a:rPr"))
            return copy.deepcopy(rpr) if rpr is not None else None
    return None


def _new_run(paragraph, texte, base_rpr, couleur=None):
    """Ajoute un run en appliquant le format cloné (base_rpr)."""
    run = paragraph.add_run()
    run.text = str(texte)
    if base_rpr is not None:
        existant = run._r.find(qn("a:rPr"))
        if existant is not None:
            run._r.remove(existant)
        run._r.insert(0, copy.deepcopy(base_rpr))
    if couleur is not None:
        run.font.color.rgb = couleur  # écrase la couleur (appuis colorés)
    return run


def _set_cell(cell, lignes, couleur=None):
    """Réécrit le texte d'une cellule (une ligne = un paragraphe), style préservé."""
    tf = cell.text_frame
    base = _base_rpr(tf)
    tf.clear()
    for i, ligne in enumerate(lignes):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _new_run(p, ligne, base, couleur)


def _remplacer_texte_global(prs, mapping):
    """Remplace des chaînes exactes au niveau des runs (préserve tout le style)."""
    for slide in prs.slides:
        for sh in slide.shapes:
            frames = []
            if sh.has_text_frame:
                frames.append(sh.text_frame)
            if sh.has_table:
                for row in sh.table.rows:
                    for c in row.cells:
                        frames.append(c.text_frame)
            for tf in frames:
                for p in tf.paragraphs:
                    for r in p.runs:
                        for old, new in mapping.items():
                            if old and old in r.text:
                                r.text = r.text.replace(old, new)


def _remplacer_image(slide, pic, chemin):
    """Remplace le blob d'une image en conservant sa position/taille."""
    try:
        image_part, rId = slide.part.get_or_add_image_part(chemin)
        blip = pic._element.find(".//" + qn("a:blip"))
        if blip is not None:
            blip.set(qn("r:embed"), rId)
            return True
    except Exception as e:
        logger.warning(f"APD PPTX : remplacement image impossible ({e})")
    return False


def _supprimer_shape(shape):
    el = shape._element
    el.getparent().remove(el)


def _supprimer_slide(prs, slide):
    """Supprime une slide de la présentation."""
    sldId_lst = prs.slides._sldIdLst
    rId = None
    for sldId in list(sldId_lst):
        if sldId.get(qn("r:id")):
            # retrouver la slide correspondante
            pass
    # méthode robuste : par index
    for i, s in enumerate(prs.slides):
        if s == slide:
            xml_slides = list(sldId_lst)
            sldId_lst.remove(xml_slides[i])
            return True
    return False


def _images_zone(slide, lmin, tmin, tmax):
    """Pictures situées dans une zone (grande image de contenu)."""
    out = []
    for sh in slide.shapes:
        if sh.shape_type == 13 and sh.top is not None:  # PICTURE
            if sh.top > Inches(tmin) and sh.top < Inches(tmax) and sh.width and sh.width > Inches(3):
                out.append(sh)
    return out


# ---------------------------------------------------------------------------
# Données métier
# ---------------------------------------------------------------------------

def _lire(dossier, nom):
    p = os.path.join(dossier, f"{nom}.shp")
    if not os.path.exists(p):
        return None
    try:
        g = gpd.read_file(p)
        return g if len(g) else None
    except Exception:
        return None


def _v(row, champ):
    x = row.get(champ)
    return "" if x is None or (isinstance(x, float) and math.isnan(x)) else str(x).strip()


def _boites_vignettes(dossier_shape, site="NRO"):
    """Vignettes BPE (nom, adresse, chambre/appui, modèle).

    Règle métier : site **BTS** -> uniquement les BPE **À CRÉER** (hors EN
    SERVICE) ; site **NRO/NRA** -> **TOUTES** les BPE (existantes + à créer)."""
    bpe = _lire(dossier_shape, "BPE")
    pt = _lire(dossier_shape, "PT")
    out = []
    if bpe is None:
        return out
    bts_only = (site == "BTS")
    for _, r in bpe.iterrows():
        if bts_only and "SERVICE" in _v(r, "ETAT").upper():
            continue  # BTS : boîtes existantes exclues
        g = r.geometry
        ptref = ""
        if g is not None and pt is not None:
            try:
                pr = pt.iloc[int(pt.distance(g).idxmin())]
                ptref = (_v(pr, "REF_CHAMBR") + " " + _v(pr, "NOM")).strip() or _v(pr, "NOM")
            except Exception:
                ptref = ""
        modele = _v(r, "MODELE")
        ref = _v(r, "REFERENCE")
        type_boite = f"{modele} - {ref}".strip(" -")
        out.append({
            "nom": _v(r, "NOM") or _v(r, "CODE"),
            "adresse": _v(r, "ADRESSE"),
            "cp_ville": f"{_v(r, 'CP')} {_v(r, 'VILLE')}".strip(),
            "ptref": ptref,
            "type_boite": type_boite,
        })
    return out


# ---------------------------------------------------------------------------
# Remplissage des slides
# ---------------------------------------------------------------------------

def _fmt_ml(v):
    try:
        return f"{int(round(float(v)))}"
    except (TypeError, ValueError):
        return "0"


def _remplir_synthese(shape, d):
    """Réécrit le cadre « Caractéristiques de la liaison » depuis les données Console."""
    s = d.get("souterrain", {})
    a = d.get("aerien", {})
    ap = d.get("appuis", {})
    b = d.get("boites", {})
    tf = shape.text_frame
    base = _base_rpr(tf)
    # Le 1er run cloné est le titre SOULIGNÉ : on retire ce soulignement du gabarit
    # pour que les lignes de données ne soient pas soulignées (sinon les
    # tabulations apparaissent comme des « traits »). Il est réappliqué
    # explicitement aux seuls intitulés qui le nécessitent (souligne=True).
    if base is not None:
        base.attrib.pop("u", None)
    tf.clear()

    def ligne(txt="", couleur=None, gras=None, souligne=False):
        p = tf.paragraphs[0] if not getattr(ligne, "_init", False) else tf.add_paragraph()
        ligne._init = True
        run = _new_run(p, txt, base, couleur)
        run.font.size = Pt(13)   # taille homogène pour tenir au-dessus du pied
        if gras is not None:
            run.font.bold = gras
        if souligne:
            run.font.underline = True
        return p
    ligne._init = False

    vtl_vis = s.get("vtl_visible", True)
    total_sout = (int(round(float(s.get("blo_ml") or 0))) + int(round(float(s.get("free_ml") or 0)))
                  + int(round(float(s.get("tiers_ml") or 0))) + int(round(float(s.get("gc_ml") or 0)))
                  + (int(round(float(s.get("vtl_ml") or 0))) if vtl_vis else 0))

    ligne("Caractéristiques de la liaison :", gras=True, souligne=True)
    ligne("")
    ligne("Réseaux Souterrain")
    ligne(f"\tBLO : {_fmt_ml(s.get('blo_ml'))}ml + {_fmt_ml(s.get('blo_ch'))} chambre(s) traversée(s)")
    ligne(f"\tRéseaux FREE : {_fmt_ml(s.get('free_ml'))}ml + {_fmt_ml(s.get('free_ch'))} chambre(s) traversée(s)")
    tn = s.get("tiers_nom") or "préciser"
    ligne(f"\tRéseaux Tiers ({tn}) : {_fmt_ml(s.get('tiers_ml'))}ml + {_fmt_ml(s.get('tiers_ch'))} chambre(s) traversée(s)")
    ligne(f"\tGC : {_fmt_ml(s.get('gc_ml'))}ml + {_fmt_ml(s.get('gc_ch'))} chambre(s) traversée(s)")
    if vtl_vis:
        ligne(f"\tVTL : {_fmt_ml(s.get('vtl_ml'))}ml")
    ligne(f"\tTOTAL réseaux souterrain : {total_sout}ml")
    ligne(f"\tTOTAL réseaux souterrain + love manchon : {_fmt_ml(s.get('total_love_ml'))}ml")
    ligne("")
    ligne("Réseaux Aériens")
    ligne(f"\tAERIEN FREE : {_fmt_ml(a.get('free_ml'))}ml + {_fmt_ml(a.get('free_ap'))} appui(s)")
    ligne(f"\tAERIEN ORANGE : {_fmt_ml(a.get('orange_ml'))}ml + {_fmt_ml(a.get('orange_ap'))} appui(s)")
    ligne(f"\tAERIEN ENEDIS : {_fmt_ml(a.get('enedis_ml'))}ml + {_fmt_ml(a.get('enedis_ap'))} appui(s)")
    tna = a.get("tiers_nom") or "préciser"
    ligne(f"\tAERIEN TIERS ({tna}) : {_fmt_ml(a.get('tiers_ml'))}ml + {_fmt_ml(a.get('tiers_ap'))} appui(s)")
    ligne(f"\tTOTAL réseaux aériens : {_fmt_ml(a.get('total_ml'))}ml + {_fmt_ml(a.get('total_ap'))} appui(s)")
    ligne(f"\tTOTAL réseaux aériens + love manchon : {_fmt_ml(a.get('total_love_ml'))}ml")
    ligne("")
    ligne(f"\tNombre d'appuis à remplacer : {_fmt_ml(ap.get('remplacer'))}", couleur=_ROUGE)
    ligne(f"\tNombre d'appuis à recaler/renforcer : {_fmt_ml(ap.get('recaler'))}", couleur=_ORANGE)
    ligne(f"\tNombre d'appuis à implanter : {_fmt_ml(ap.get('implanter'))}", couleur=_BLEU)
    ligne("")
    ligne("Boites")

    def bo(v):
        return _fmt_ml(v) if (v not in (None, "", 0)) else ""
    ligne(f"\tNombre de boite FREE OFDC : {bo(b.get('ofdc'))}")
    ligne(f"\tNombre de boite FREE TENIO : {bo(b.get('tenio'))}")
    ligne(f"\tNombre de boite FREE 3M T0 : {bo(b.get('t0'))}")
    ligne(f"\tNombre de boite FREE 3M T1 : {bo(b.get('t1'))}")
    ligne(f"\tTOTAL boites : {_fmt_ml(b.get('total'))}")
    ligne("")
    ligne("Informations remarquables :", gras=True, souligne=True)
    infos = (d.get("infos") or "").strip()
    if infos:
        for l in infos.split("\n"):
            ligne(l)


def _remplir_footer(slide, cart):
    ft = _footer_table(slide)
    if ft is None:
        return
    t = ft.table
    valeurs = [
        ["Version du document :", cart.get("version", "")],
        ["Code projet du site :", cart.get("code_projet", "")],
        ["Fait par :", cart.get("fait_par", "")],
        ["Date de réalisation :", ("Le " + cart.get("date_real", "")) if cart.get("date_real") else ""],
    ]
    for i, cell in enumerate(t.rows[0].cells):
        if i < len(valeurs):
            _set_cell(cell, valeurs[i])


def _remplir_header(slide, cart):
    ht = _header_table(slide)
    if ht is None:
        return
    t = ht.table
    # cellule centrale = code + type
    site = cart.get("type_etude", "APD HTL").replace("APD HTL", "").strip() or ""
    _set_cell(t.rows[0].cells[1], [cart.get("code_projet", ""), "APD HTL      " + site])


def _remplir_vignettes(slide, boites):
    """Remplace les vignettes du template par une vignette par boîte à créer."""
    vign = [sh for sh in slide.shapes if sh.has_table and sh.top is not None
            and sh.top > Inches(1) and sh.top < Inches(9.5) and len(sh.table.columns) == 1]
    vign.sort(key=lambda s: s.top)
    if not vign:
        return
    proto = vign[0]._element
    largeur = vign[0].width
    hauteur = vign[0].height or Inches(0.95)
    # supprimer les vignettes existantes
    for sh in vign:
        _supprimer_shape(sh)

    spTree = slide.shapes._spTree
    col_x = [Inches(0.41), Inches(4.25)]
    top0, step = 1.13, 1.25
    par_col = max(1, int((10.4 - top0) / step))

    for i, box in enumerate(boites):
        el = copy.deepcopy(proto)
        spTree.append(el)
        shape = slide.shapes[-1]
        col = i // par_col
        row = i % par_col
        shape.left = col_x[min(col, 1)]
        shape.top = Inches(top0 + row * step)
        shape.width = largeur
        t = shape.table
        # r0 = nom ; r1 = adresse + cp/ville ; r2 = chambre/appui + type boîte
        _set_cell(t.rows[0].cells[0], [box["nom"]])
        _set_cell(t.rows[1].cells[0], [f"Adresse : {box['adresse']}", box["cp_ville"]])
        _set_cell(t.rows[2].cells[0],
                  [f"N° de la chambre/appui : {box['ptref']}", f"Type de boîte : {box['type_boite']}"])


# ---------------------------------------------------------------------------
# Orchestration
# ---------------------------------------------------------------------------

def remplir_apd(template_path, donnees, dossier_shape, sortie_pptx,
                plan_general=None, plan_masse=None):
    """Copie le template et le remplit ; renvoie le chemin du PPTX."""
    prs = Presentation(template_path)
    cart = donnees.get("cartouche", {})
    site = "BTS" if "BTS" in cart.get("type_etude", "") else (
        "NRO" if "NRO" in cart.get("type_etude", "") else "NRA")

    slides = list(prs.slides)
    # En-tête + pied sur toutes les slides
    for s in slides:
        _remplir_header(s, cart)
        _remplir_footer(s, cart)

    # --- SLIDE 1 : adresse + plan général ---
    if len(slides) >= 1:
        s1 = slides[0]
        # adresse = la "Forme automatique" qui n'est pas "PLAN"
        for sh in s1.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip() and sh.text_frame.text.strip().upper() != "PLAN":
                if sh.left is not None and sh.left < Inches(8):  # sur le canvas
                    ad = cart.get("adresse", "")
                    ligne_adr = ad
                    if cart.get("cp") or cart.get("commune"):
                        ligne_adr = f"{ad} - {cart.get('cp','')} - {cart.get('commune','')}".strip(" -")
                    _set_cell_like(sh, ligne_adr)
                    break
        if plan_general and os.path.exists(plan_general):
            imgs = _images_zone(s1, 0.41, 2.0, 9.0)
            for j, pic in enumerate(imgs):
                if j == 0:
                    _remplacer_image(s1, pic, plan_general)
                else:
                    _supprimer_shape(pic)

    # --- SLIDE 2 : caractéristiques de la liaison ---
    if len(slides) >= 2:
        s2 = slides[1]
        cible = None
        for sh in s2.shapes:
            if sh.has_text_frame and sh.left is not None and sh.left < Inches(8):
                if "aract" in sh.text_frame.text and "liaison" in sh.text_frame.text:
                    cible = sh
                    break
        if cible is not None:
            _remplir_synthese(cible, donnees)

    # --- SLIDE 3 : plan de masse — gardée UNIQUEMENT si BTS + image fournie ---
    if len(slides) >= 3:
        s3 = slides[2]
        a_plan = bool(plan_masse and os.path.exists(plan_masse))
        if site != "BTS" or not a_plan:
            _supprimer_slide(prs, s3)
        else:
            imgs = _images_zone(s3, 0.71, 1.0, 10.0)
            for j, pic in enumerate(imgs):
                if j == 0:
                    _remplacer_image(s3, pic, plan_masse)
                else:
                    _supprimer_shape(pic)

    # --- SLIDE 4 : vignettes BPE à créer ---
    s4 = None
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_text_frame and "vignette" in sh.text_frame.text.lower():
                s4 = s
                break
    if s4 is None and len(list(prs.slides)) >= 4:
        s4 = list(prs.slides)[-1]
    if s4 is not None:
        boites = _boites_vignettes(dossier_shape, site)
        if boites:
            _remplir_vignettes(s4, boites)      # remplace les vignettes-exemple du template
        else:
            _supprimer_slide(prs, s4)           # aucune boîte -> pas de page vignettes-exemple

    os.makedirs(os.path.dirname(sortie_pptx), exist_ok=True)
    prs.save(sortie_pptx)
    logger.info(f"APD PPTX rempli : {sortie_pptx}")
    return sortie_pptx


def _set_cell_like(shape, texte):
    """Réécrit un cadre texte simple en préservant sa police."""
    tf = shape.text_frame
    base = _base_rpr(tf)
    tf.clear()
    _new_run(tf.paragraphs[0], texte, base)


def pptx_vers_pdf(pptx_path, pdf_path):
    """Export PPTX -> PDF via PowerPoint COM, dans un SOUS-PROCESSUS isolé
    (instance PowerPoint fraîche par appel — robuste face aux threads serveur)."""
    import sys
    import subprocess
    script = os.path.join(os.path.dirname(__file__), "pptx_to_pdf.py")
    if os.path.exists(pdf_path):
        try:
            os.remove(pdf_path)
        except Exception:
            pass
    res = subprocess.run(
        [sys.executable, script, os.path.abspath(pptx_path), os.path.abspath(pdf_path)],
        capture_output=True, text=True, timeout=180,
    )
    if res.returncode != 0 or not os.path.exists(pdf_path):
        raise RuntimeError(
            "Export PPTX→PDF échoué (PowerPoint requis). "
            f"code={res.returncode} : {(res.stderr or res.stdout or '').strip()[-300:]}"
        )
    return pdf_path


def generer_apd_plan(template_path, donnees, dossier_shape, sortie_pdf,
                     plan_general=None, plan_masse=None):
    """Remplit le template puis exporte en PDF. Renvoie le chemin du PDF."""
    sortie_pptx = os.path.splitext(sortie_pdf)[0] + ".pptx"
    remplir_apd(template_path, donnees, dossier_shape, sortie_pptx,
                plan_general=plan_general, plan_masse=plan_masse)
    pptx_vers_pdf(sortie_pptx, sortie_pdf)
    logger.info(f"APD plan PDF généré : {sortie_pdf}")
    return sortie_pdf
